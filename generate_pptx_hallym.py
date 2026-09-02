from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
# 16:9 와이드스크린 슬라이드 크기 설정
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6]) # 빈 슬라이드

# 스타일 도우미 함수
def add_card(slide, left, top, width, height, title, body_text, title_bg, body_bg):
    # 배경 박스
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = body_bg
    rect.line.color.rgb = title_bg
    rect.line.width = Pt(1.5)
    
    # 타이틀 바
    t_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.55))
    t_bar.fill.solid()
    t_bar.fill.fore_color.rgb = title_bg
    t_bar.line.fill.background()
    tf_t = t_bar.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title
    p_t.font.name = "Malgun Gothic"
    p_t.font.size = Pt(11)
    p_t.font.bold = True
    p_t.font.color.rgb = RGBColor(255, 255, 255)
    p_t.alignment = PP_ALIGN.CENTER

    # 본문 텍스트 박스
    if body_text:
        tx_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.6), width - Inches(0.2), height - Inches(0.7))
        tf_b = tx_box.text_frame
        tf_b.word_wrap = True
        for i, line in enumerate(body_text.strip().split('\n')):
            p = tf_b.paragraphs[0] if i == 0 else tf_b.add_paragraph()
            p.text = line
            p.font.name = "Malgun Gothic"
            p.font.size = Pt(9.5)
            p.font.color.rgb = RGBColor(33, 37, 41)

# --- 1-Tier: Client ---
add_card(slide, Inches(0.6), Inches(0.6), Inches(5.8), Inches(2.5), 
         "1-Tier : Presentation Tier (원내 PC 단말)", "", RGBColor(43, 108, 176), RGBColor(240, 244, 248))

# 1-Tier 내부 박스들
add_card(slide, Inches(0.9), Inches(1.3), Inches(2.4), Inches(1.5), 
         "원내 EMR / CPOE", "• 처방 입력 화면\n• 실시간 경고 팝업", RGBColor(74, 85, 104), RGBColor(255, 255, 255))
add_card(slide, Inches(3.6), Inches(1.3), Inches(2.4), Inches(1.5), 
         "약품정보 전용 조회", "• .NET Client App\n• 약품 상세 가이드", RGBColor(49, 130, 206), RGBColor(235, 248, 255))

# --- 2-Tier: AP Server ---
add_card(slide, Inches(6.8), Inches(0.6), Inches(5.9), Inches(2.5), 
         "2-Tier : Business Logic Tier (AP 서버)", 
         "• .NET 비즈니스 로직 및 컴포넌트\n• 약품 상세 효능·용법·주의사항 데이터 처리\n• 최신 약가 고시 및 제약 마스터 서비스 가공", 
         RGBColor(74, 85, 104), RGBColor(247, 250, 252))

# --- 3-Tier: Database Container ---
add_card(slide, Inches(0.6), Inches(3.4), Inches(12.1), Inches(3.6), 
         "3-Tier : Database Tier (원내 EMR DBMS 환경 - Oracle / MSSQL / PostgreSQL 등)", "", RGBColor(26, 54, 93), RGBColor(237, 242, 247))

# 3-Tier 내부 (EMR DB vs DIF DB)
add_card(slide, Inches(1.0), Inches(4.1), Inches(5.4), Inches(2.6), 
         "원내 EMR Database", 
         "★ 실시간 CDSS Stored Procedure\n• CPOE 처방 발생 시 DB 내부 직접 호출\n• 병용/연령/임부 금기 및 용량 실시간 판별\n• 네트워크 무지연 (Zero-Latency) 초고속 연산", 
         RGBColor(43, 108, 176), RGBColor(254, 252, 191))

add_card(slide, Inches(6.9), Inches(4.1), Inches(5.4), Inches(2.6), 
         "DIF Database (원내 EMR DB 내 스키마 탑재)", 
         "DIF 약품 마스터 & 임상 룰셋 데이터\n• 의약품 표준 식별 및 제약 마스터 정보\n• 복약지도, 상호작용 및 임상 기준 데이터\n• EMR 트랜잭션 동기화 및 무결성 확보", 
         RGBColor(49, 151, 149), RGBColor(230, 255, 250))

prs.save("system_architecture.pptx")
print("PPTX Created Successfully.")
